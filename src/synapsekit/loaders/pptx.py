from __future__ import annotations

import os

from .base import Document


class PowerPointLoader:
    """Load a PowerPoint (.pptx) file, one Document per slide.

    Extracts text from all shapes (text frames) on each slide.
    """

    def __init__(self, path: str) -> None:
        self._path = path

    def load(self) -> list[Document]:
        if not os.path.exists(self._path):
            raise FileNotFoundError(f"PowerPoint file not found: {self._path}")
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx required: pip install synapsekit[pptx]") from None

        prs = Presentation(self._path)
        docs = []
        for i, slide in enumerate(prs.slides):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            parts.append(text)
            docs.append(
                Document(
                    text="\n".join(parts),
                    metadata={"source": self._path, "slide": i},
                )
            )
        return docs
